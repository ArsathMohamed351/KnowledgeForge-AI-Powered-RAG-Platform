import os
import json
import pandas as pd
from langchain_core.documents import Document
from langchain_community.document_loaders import PyPDFLoader
from docx import Document as DocxDocument
from pptx import Presentation

CSV_BATCH_SIZE = 100
def create_metadata(filename, file_path, file_type):
    return { "source": filename, "file_path": file_path, "file_type": file_type }

# PDF And TXT

def load_file(file_path):
    extension = os.path.splitext(file_path)[1].lower()
    filename = os.path.basename(file_path)
    docs = []
    if extension == ".pdf":
        loader = PyPDFLoader(file_path)
        docs = loader.load()
        for doc in docs:
            doc.metadata.update(create_metadata(filename, file_path, "pdf"))
    elif extension == ".txt":
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()
        docs = [
            Document( page_content=text, metadata={ **create_metadata(filename, file_path, "txt"), "dataset_id": file_path } )
        ]
# CSV

    elif extension == ".csv":
        df = pd.read_csv(file_path, low_memory=False)
        total_rows = len(df)

        for start in range(0, total_rows, CSV_BATCH_SIZE):
            batch_df = df.iloc[start:start + CSV_BATCH_SIZE]
            docs.append(
                Document(
                    page_content=batch_df.to_string(index=False),
                    metadata={
                        **create_metadata(filename, file_path, "csv"),
                        "dataset_id": file_path,
                        "batch_start": start,
                        "batch_end": min(start + CSV_BATCH_SIZE, total_rows)
                    }
                )
            )

# Excel Format (xlsx, xls)
    elif extension in [".xlsx", ".xls"]:
        df = pd.read_excel(file_path)
        total_rows = len(df)
        for start in range(0, total_rows, CSV_BATCH_SIZE):
            batch_df = df.iloc[start:start + CSV_BATCH_SIZE]
            docs.append(
                Document( page_content=batch_df.to_string(index=False), metadata={  **create_metadata(filename, file_path, "xlsx"),  "dataset_id": file_path,  "batch_start": start,  "batch_end": min(start + CSV_BATCH_SIZE, total_rows)} ) )

# DOCX

    elif extension == ".docx":
        doc = DocxDocument(file_path)
        text = "\n".join(p.text for p in doc.paragraphs)
        docs = [
            Document( page_content=text, metadata=create_metadata(filename, file_path, "docx") )
        ]

# PPTX

    elif extension == ".pptx":
        presentation = Presentation(file_path)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        docs = [
            Document( page_content=text, metadata=create_metadata(filename, file_path, "pptx") )
        ]

# JSON

    elif extension == ".json":
        with open(file_path, "r", encoding="utf-8") as f:
            data = json.load(f)
        docs = [
            Document( page_content=json.dumps(data, indent=2), metadata=create_metadata(filename, file_path, "json") )
        ]
    else:
        raise ValueError(f"Unsupported file type: {extension}")
    return docs